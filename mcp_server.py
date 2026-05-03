from fastmcp import FastMCP
import os
import subprocess
import re
import json
import shutil
import time
import difflib
import plotly.express as px
import pandas as pd
from graphviz import Digraph
from typing import Optional
from duckduckgo_search import DDGS
import requests
from bs4 import BeautifulSoup
import fnmatch
import html as html_lib
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor

# Optional document libraries
try:
    from docx import Document
    from docx.shared import Inches as DocxInches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    Document = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PptInches, Pt as PptPt
except ImportError:
    Presentation = None

try:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Image as RLImage,
        Table as RLTable, TableStyle, PageBreak
    )
    from reportlab.lib.utils import ImageReader
except ImportError:
    SimpleDocTemplate = None

# RAG knowledge base bridge (optional — requires RAG_PROJECT_PATH in .env)
try:
    from rag_tool_bridge import (
        rag_search,
        rag_semantic_search,
        rag_search_by_book,
        rag_search_by_section,
        rag_health_check,
    )
except ImportError:
    def rag_search(*a, **kw): return "RAG bridge not available"
    def rag_semantic_search(*a, **kw): return "RAG bridge not available"
    def rag_search_by_book(*a, **kw): return "RAG bridge not available"
    def rag_search_by_section(*a, **kw): return "RAG bridge not available"
    def rag_health_check(*a, **kw): return "RAG bridge not available"


mcp = FastMCP("react-agent-tools")

BASE_DIR = os.path.realpath(os.getcwd())
BACKUP_DIR = os.path.join(BASE_DIR, ".backups")


def safe_path(path: str) -> str:
    full = os.path.realpath(os.path.join(BASE_DIR, path))
    if not (full == BASE_DIR or full.startswith(BASE_DIR + os.sep)):
        raise ValueError("Access outside sandbox is not allowed")
    return full


def _backup_file(path: str) -> Optional[str]:
    """Create a timestamped backup of a file before modification. Returns backup path or None."""
    if not os.path.isfile(path):
        return None
    os.makedirs(BACKUP_DIR, exist_ok=True)
    rel = os.path.relpath(path, BASE_DIR)
    safe_name = rel.replace(os.sep, "__")
    ts = int(time.time())
    backup_path = os.path.join(BACKUP_DIR, f"{safe_name}.{ts}.bak")
    shutil.copy2(path, backup_path)
    return backup_path


def _get_diff(original_path: str, new_content: str, context_lines: int = 3) -> str:
    """Generate a unified diff between original file and new content."""
    try:
        with open(original_path, "r", encoding="utf-8", errors="replace") as f:
            original_lines = f.readlines()
        new_lines = new_content.splitlines(keepends=True)
        diff = difflib.unified_diff(
            original_lines, new_lines,
            fromfile="original", tofile="modified",
            n=context_lines,
        )
        return "".join(diff)
    except Exception:
        return "(diff unavailable)"


# ───────────────── FILE TOOLS ───────────────── #

@mcp.tool()
def list_files(path: str = ".") -> str:
    """List directory contents with [DIR]/[FILE] indicators."""
    try:
        path = safe_path(path)
    except ValueError as e:
        return str(e)

    if not os.path.exists(path):
        return f"Path does not exist: {path}"

    if not os.path.isdir(path):
        return f"Not a directory: {path}"

    items = []
    for item in sorted(os.listdir(path)):
        full = os.path.join(path, item)
        if os.path.isdir(full):
            items.append(f"[DIR]  {item}/")
        else:
            size = os.path.getsize(full)
            items.append(f"[FILE] {item}  ({size} bytes)")
    return "\n".join(items) if items else "(empty directory)"


@mcp.tool()
def read_file(path: str, start_line: Optional[int] = None, end_line: Optional[int] = None) -> str:
    """Read a file with line numbers. Optionally specify start_line and end_line (1-indexed)."""
    try:
        path = safe_path(path)
    except ValueError as e:
        return str(e)

    if not os.path.exists(path):
        return f"File does not exist: {path}"

    if not os.path.isfile(path):
        return f"Not a file: {path}"

    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            lines = f.readlines()
    except PermissionError:
        return f"Permission denied: {path}"

    total = len(lines)

    start = (start_line - 1) if start_line else 0
    end = end_line if end_line else total

    start = max(0, start)
    end = min(total, end)

    selected = lines[start:end]

    result = []
    for i, line in enumerate(selected):
        result.append(f"{start + i + 1:4d} | {line.rstrip()}")

    return f"{path} ({total} lines)\n" + "\n".join(result)


@mcp.tool()
def write_file(path: str, content: str) -> str:
    """Write content to a file. Creates parent directories if needed. Auto-backs up existing files."""
    try:
        path = safe_path(path)
    except ValueError as e:
        return str(e)

    if os.path.isfile(path):
        _backup_file(path)

    parent_dir = os.path.dirname(path)
    if parent_dir:
        os.makedirs(parent_dir, exist_ok=True)

    try:
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
    except PermissionError:
        return f"Permission denied: {path}"
    except OSError as e:
        return f"Write error: {e}"

    line_count = content.count("\n") + (1 if content and not content.endswith("\n") else 0)
    return f"Written to {path} ({line_count} lines, {len(content)} bytes)"


@mcp.tool()
def patch_file(path: str, search: str, replace: str) -> str:
    """Surgical search-and-replace in a file. Replaces the first occurrence of `search` with `replace`. Shows diff."""
    try:
        path = safe_path(path)
    except ValueError as e:
        return str(e)

    if not os.path.exists(path):
        return f"File does not exist: {path}"

    try:
        with open(path, "r", encoding="utf-8") as f:
            data = f.read()
    except PermissionError:
        return f"Permission denied: {path}"

    if search not in data:
        return f"Search text not found in {path}"

    new_data = data.replace(search, replace, 1)

    _backup_file(path)
    diff = _get_diff(path, new_data)

    try:
        with open(path, "w", encoding="utf-8") as f:
            f.write(new_data)
    except PermissionError:
        return f"Permission denied: {path}"
    except OSError as e:
        return f"Write error: {e}"

    return f"Patch applied to {path}\n--- diff ---\n{diff}"


@mcp.tool()
def delete_file(path: str) -> str:
    """Delete a file (auto-backs up first)."""
    try:
        path = safe_path(path)
    except ValueError as e:
        return str(e)

    if not os.path.exists(path):
        return f"File does not exist: {path}"

    if not os.path.isfile(path):
        return f"Not a file: {path}"

    _backup_file(path)
    try:
        os.remove(path)
        return f"Deleted: {path}"
    except PermissionError:
        return f"Permission denied: {path}"
    except OSError as e:
        return f"Delete error: {e}"


@mcp.tool()
def revert_file(path: str) -> str:
    """Revert a file to its most recent backup. Shows diff before reverting."""
    try:
        path = safe_path(path)
    except ValueError as e:
        return str(e)

    if not os.path.isdir(BACKUP_DIR):
        return "No backups available"

    rel = os.path.relpath(path, BASE_DIR)
    safe_name = rel.replace(os.sep, "__")
    prefix = f"{safe_name}."

    backups = []
    for f in os.listdir(BACKUP_DIR):
        if f.startswith(prefix) and f.endswith(".bak"):
            # Extract timestamp: filename is "{safe_name}.{timestamp}.bak"
            middle = f[len(prefix):-4]  # strip prefix and ".bak"
            try:
                ts = int(middle)
                backups.append((ts, f))
            except ValueError:
                continue

    if not backups:
        return f"No backups found for: {path}"

    backups.sort(key=lambda x: x[0], reverse=True)
    latest_ts, latest_name = backups[0]
    latest_backup = os.path.join(BACKUP_DIR, latest_name)

    if os.path.isfile(path):
        _backup_file(path)
        with open(latest_backup, "r", encoding="utf-8", errors="replace") as f:
            backup_content = f.read()
        diff = _get_diff(path, backup_content)
    else:
        diff = "(file will be restored from backup)"

    try:
        shutil.copy2(latest_backup, path)
        return f"Reverted {path} to backup from {time.ctime(latest_ts)}\n--- diff ---\n{diff}"
    except Exception as e:
        return f"Revert error: {str(e)}"


@mcp.tool()
def list_backups() -> str:
    """List all file backups with timestamps."""
    if not os.path.isdir(BACKUP_DIR):
        return "No backups available"

    backups = []
    for f in sorted(os.listdir(BACKUP_DIR)):
        if f.endswith(".bak"):
            # Parse: "{name}.{timestamp}.bak"
            without_bak = f[:-4]  # strip ".bak"
            dot_pos = without_bak.rfind(".")
            if dot_pos > 0:
                name = without_bak[:dot_pos]
                ts_str = without_bak[dot_pos + 1:]
                try:
                    ts = int(ts_str)
                    original_name = name.replace("__", os.sep)
                    t = time.ctime(ts)
                    size = os.path.getsize(os.path.join(BACKUP_DIR, f))
                    backups.append(f"{t}  {original_name}  ({size} bytes)")
                except ValueError:
                    continue

    if not backups:
        return "No backups available"

    return "\n".join(backups)


# ───────────────── SEARCH TOOLS ───────────────── #

@mcp.tool()
def grep_search(pattern: str, path: str = ".", include: Optional[str] = None, max_results: int = 50) -> str:
    """Search for a text pattern in files under `path`. Returns matching lines with file:line context.

    Args:
        pattern: Text pattern to search for (case-insensitive)
        path: Directory or file to search in (default: current directory)
        include: Optional glob to filter filenames (e.g. '*.py')
        max_results: Maximum number of matches to return (default: 50)
    """
    try:
        search_path = safe_path(path)
    except ValueError as e:
        return str(e)

    if not os.path.exists(search_path):
        return f"Path does not exist: {search_path}"

    matches = []
    pattern_lower = pattern.lower()

    def search_file(filepath: str):
        try:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                for line_num, line in enumerate(f, 1):
                    if pattern_lower in line.lower():
                        rel = os.path.relpath(filepath, BASE_DIR)
                        matches.append(f"{rel}:{line_num}: {line.rstrip()}")
                        if len(matches) >= max_results:
                            return
        except (PermissionError, IsADirectoryError, OSError):
            pass

    if os.path.isfile(search_path):
        search_file(search_path)
    else:
        for root, dirs, files in os.walk(search_path):
            # Skip hidden/generated directories
            dirs[:] = [d for d in dirs if not d.startswith(".") and d not in ("__pycache__", "node_modules", ".venv", ".git")]
            for fname in files:
                if include and not fnmatch.fnmatch(fname, include):
                    continue
                filepath = os.path.join(root, fname)
                search_file(filepath)
                if len(matches) >= max_results:
                    break
            if len(matches) >= max_results:
                break

    if not matches:
        return f"No matches found for '{pattern}' in {path}"

    header = f"Found {len(matches)} match(es) for '{pattern}':\n"
    return header + "\n".join(matches)


@mcp.tool()
def glob_search(pattern: str, path: str = ".") -> str:
    """Find files matching a glob pattern (e.g. '*.py', '**/*.js', 'src/**/*.ts').

    Args:
        pattern: Glob pattern to match filenames against
        path: Directory to search in (default: current directory)
    """
    try:
        search_path = safe_path(path)
    except ValueError as e:
        return str(e)

    if not os.path.exists(search_path):
        return f"Path does not exist: {search_path}"

    import pathlib
    base = pathlib.Path(search_path)
    results = []
    for match in base.glob(pattern):
        # Skip hidden dirs and common noise
        parts = match.relative_to(base).parts
        if any(p.startswith(".") or p in ("__pycache__", "node_modules", ".venv") for p in parts):
            continue
        rel = os.path.relpath(str(match), BASE_DIR)
        kind = "[DIR]" if match.is_dir() else "[FILE]"
        results.append(f"{kind} {rel}")
        if len(results) >= 200:
            results.append(f"... (truncated at 200 results)")
            break

    if not results:
        return f"No files matching '{pattern}' in {path}"

    return f"Found {len(results)} match(es):\n" + "\n".join(results)


# ───────────────── BASH ───────────────── #

# Block shell metacharacters that enable command injection, chaining, or
# file-system manipulation via redirection.
# Covers: ; | & ` (chaining/subshell), $( ) (command substitution),
#         > >> < (redirection), ! (history expansion).
_SHELL_METACHARACTERS = re.compile(r'[;|&`><!]|\$\(|\|\||&&')

DANGEROUS_COMMANDS = [
    r"\brm\s+(-rf?|--recursive)\s+/$",
    r"\brm\s+(-rf?|--recursive)\s+\*",
    r"\bformat\b",
    r"\bdel\s+C:\\",
    r"\bmkfs\b",
    r"\bdd\s+if=/dev/zero\b",
    r"\bshutdown\b",
    r"\breboot\b",
    r"\bpkill\s+-9\b",
    r"\bkillall\b",
    r"\bsudo\b",
    r"\bchmod\s+[0-7]*777\b",
    r"\bchown\b",
]

ALLOWED_COMMANDS = [
    "ls", "cat", "grep", "find", "pwd", "echo", "head", "tail",
    "wc", "sort", "uniq", "diff", "tree", "file", "stat", "whoami",
    "git", "python", "python3", "pip", "pip3", "npm", "node",
    "make", "cmake", "gcc", "g++", "cargo", "rustc",
    "curl", "wget",
    "docker", "docker-compose",
    "pytest", "unittest",
    "tar", "zip", "unzip", "gzip",
    "mkdir", "touch", "cp", "mv", "rmdir", "rm",
    "awk", "sed", "cut", "tr",
    "man", "which", "whereis",
    "env",
    "ps", "top", "htop",
    "date", "uptime", "cal",
    "du", "df", "free",
    "cd",
    "dir", "type", "where",  # Windows equivalents
]


def is_command_safe(command: str) -> tuple:
    """Check if a shell command is safe to execute. Returns (is_safe, reason)."""
    cmd_stripped = command.strip()
    if not cmd_stripped:
        return False, "Empty command"

    base_cmd = cmd_stripped.split()[0].lower()
    base_cmd = os.path.basename(base_cmd)
    if base_cmd.endswith(".exe"):
        base_cmd = base_cmd[:-4]

    if base_cmd not in ALLOWED_COMMANDS:
        return False, f"Command '{base_cmd}' is not in the allowed list"

    # For python/python3 with -c flag, allow metacharacters inside the script string
    is_python = base_cmd in ("python", "python3")
    if is_python and re.search(r'\s-c\s', cmd_stripped):
        # Only check dangerous patterns, skip metacharacter check for inline scripts
        for pattern in DANGEROUS_COMMANDS:
            if re.search(pattern, cmd_stripped):
                return False, "Dangerous command pattern detected"
        return True, ""

    # Block shell metacharacters that enable command chaining/injection/redirection
    if _SHELL_METACHARACTERS.search(cmd_stripped):
        return False, "Shell metacharacters (;, |, &, `, $(), >, <, !) are not allowed. Use simple single commands only."

    if base_cmd == "rm":
        if re.search(r'\b(-r|-R|--recursive)\b', cmd_stripped):
            return False, "Recursive rm is not allowed"
        if "*" in cmd_stripped:
            return False, "Glob patterns in rm are not allowed"

    for pattern in DANGEROUS_COMMANDS:
        if re.search(pattern, cmd_stripped):
            return False, "Dangerous command pattern detected"

    return True, ""


@mcp.tool()
def bash(command: str) -> str:
    """Execute a shell command. Simple commands and python3 -c inline scripts are supported."""
    safe, reason = is_command_safe(command)
    if not safe:
        return f"Command blocked: {reason}"

    try:
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=120,
            cwd=BASE_DIR,
        )

        output = []
        if result.stdout:
            stdout = result.stdout
            if len(stdout) > 10000:
                stdout = stdout[:10000] + f"\n... (truncated, {len(result.stdout)} chars total)"
            output.append(stdout)
        if result.stderr:
            stderr = result.stderr
            if len(stderr) > 5000:
                stderr = stderr[:5000] + f"\n... (truncated)"
            output.append(f"STDERR:\n{stderr}")
        if result.returncode != 0:
            output.append(f"Exit code: {result.returncode}")

        return "\n".join(output) if output else "Command completed with no output"
    except subprocess.TimeoutExpired:
        return "Command timed out after 120 seconds"
    except Exception as e:
        return f"Execution error: {e}"


# ───────────────── CHART TOOLS ───────────────── #

def _generate_text_bar_chart(data: str, title: str = "Chart", chart_type: str = "bar") -> str:
    """Generate a text-based bar chart for terminal display."""
    try:
        rows = [list(map(float, r.split(","))) for r in data.split(";")]
        if not rows or not rows[0]:
            return ""
        df = pd.DataFrame(rows)
        if df.shape[1] < 2:
            return ""

        x_vals = [str(int(v)) if v == int(v) else str(v) for v in df[0]]
        y_vals = df[1].tolist()

        max_val = max(y_vals) if y_vals else 0
        max_bar_width = 50
        max_label_width = max(len(x) for x in x_vals) if x_vals else 4

        lines = []
        lines.append(f"\n{'=' * (max_label_width + max_bar_width + 20)}")
        lines.append(f"  {title}")
        lines.append(f"{'=' * (max_label_width + max_bar_width + 20)}")
        lines.append("")

        for x, y in zip(x_vals, y_vals):
            if max_val > 0:
                bar_len = int((y / max_val) * max_bar_width)
            else:
                bar_len = 0
            bar = "█" * bar_len
            y_str = f"{y:,.1f}" if y != int(y) else f"{y:,.0f}"
            lines.append(f"  {x:>{max_label_width}} │ {bar} {y_str}")

        lines.append(f"  {' ' * max_label_width} │")
        lines.append(f"{'=' * (max_label_width + max_bar_width + 20)}")
        lines.append("")

        return "\n".join(lines)
    except Exception:
        return ""


def _generate_unique_filename(base: str, extension: str, save_dir: str) -> str:
    """Generate a unique filename based on timestamp to avoid overwriting."""
    ts = int(time.time() * 1000)
    safe_base = re.sub(r'[^a-zA-Z0-9_-]', '_', base)[:60]
    return os.path.join(save_dir, f"{safe_base}_{ts}.{extension}")


@mcp.tool()
def generate_plotly_chart(data: str, chart_type: str = "auto", title: str = "Chart", colors: Optional[str] = None) -> str:
    """Generate a chart from CSV-like data.

    Args:
        data: Semicolon-separated rows, comma-separated values (e.g. '10,20;30,40')
        chart_type: bar, line, scatter, histogram, heatmap, or auto
        title: Chart title
        colors: Comma-separated color names or hex codes (e.g. 'blue,orange,#ff0000')
    """
    try:
        rows = [list(map(float, r.split(","))) for r in data.split(";")]
        if not rows or not rows[0]:
            return "Error: Empty data. Expected format: '1,2,3;4,5,6'"
        df = pd.DataFrame(rows)

        if chart_type == "auto":
            if df.shape[1] == 1:
                chart_type = "histogram"
            elif df.shape[1] == 2:
                chart_type = "scatter"
            else:
                chart_type = "heatmap"

        color_list = None
        if colors:
            color_list = [c.strip() for c in colors.split(",")]

        if chart_type == "bar":
            if df.shape[1] == 2:
                x_vals = [str(int(v)) if v == int(v) else str(v) for v in df[0]]
                y_vals = df[1].tolist()
                bar_color = color_list[0] if color_list else "#007AFF"
                import plotly.graph_objects as go
                fig = go.Figure(go.Bar(
                    x=x_vals,
                    y=y_vals,
                    marker_color=bar_color,
                    text=[f"{v}" for v in y_vals],
                    textposition="outside",
                ))
            else:
                fig = px.bar(df, color_discrete_sequence=color_list) if color_list else px.bar(df)
        elif chart_type == "line":
            if df.shape[1] == 2:
                x_vals = [str(int(v)) if v == int(v) else str(v) for v in df[0]]
                y_vals = df[1].tolist()
                line_color = color_list[0] if color_list else "#007AFF"
                import plotly.graph_objects as go
                fig = go.Figure(go.Scatter(x=x_vals, y=y_vals, mode="lines+markers", line=dict(color=line_color, width=3)))
            else:
                fig = px.line(df, color_discrete_sequence=color_list) if color_list else px.line(df)
        elif chart_type == "scatter":
            fig = px.scatter(df, x=0, y=1, color_discrete_sequence=color_list) if color_list else px.scatter(df, x=0, y=1)
        elif chart_type == "heatmap":
            fig = px.imshow(df)
        elif chart_type == "histogram":
            fig = px.histogram(df, color_discrete_sequence=color_list) if color_list else px.histogram(df)
        else:
            return f"Unsupported chart: {chart_type}"

        fig.update_layout(
            template="plotly_dark",
            font=dict(size=14),
            title_font_size=20,
            margin=dict(l=60, r=40, t=80, b=60),
            title=dict(text=title, x=0.5),
            xaxis=dict(title="", tickfont=dict(size=13)),
            yaxis=dict(title="", tickfont=dict(size=13), gridcolor="rgba(255,255,255,0.1)"),
        )

        if chart_type == "bar":
            fig.update_traces(marker=dict(line=dict(width=0)))

        save_dir = safe_path("saves")
        os.makedirs(save_dir, exist_ok=True)

        title_slug = re.sub(r'[^a-zA-Z0-9_-]', '_', title)[:40].strip('_')
        filename = _generate_unique_filename(f"{title_slug}_{chart_type}", "png", save_dir)
        fig.write_image(filename, engine="kaleido")

        rel_path = os.path.relpath(filename, BASE_DIR)

        text_chart = _generate_text_bar_chart(data, title, chart_type)

        return json.dumps({
                        "text_chart": text_chart,
                        "path": rel_path
                    })

    except ValueError as e:
        return f"Error parsing data: {e}. Expected format: '1,2,3;4,5,6'"
    except Exception as e:
        return f"Error: {str(e)}"


@mcp.tool()
def smart_chart(user_query: str, data: str, colors: Optional[str] = None) -> str:
    """Agent-friendly chart generator. Infers chart type from query.

    Args:
        user_query: Natural language description (e.g. 'show trend', 'compare values')
        data: Semicolon-separated rows, comma-separated values
        colors: Comma-separated color names or hex codes (optional)
    """
    try:
        query = user_query.lower()

        if "trend" in query or "over time" in query:
            chart_type = "line"
        elif "compare" in query:
            chart_type = "bar"
        elif "distribution" in query:
            chart_type = "histogram"
        elif "relationship" in query:
            chart_type = "scatter"
        else:
            chart_type = "auto"

        return generate_plotly_chart(data, chart_type, user_query, colors)
    except Exception as e:
        return f"Error: {str(e)}"


@mcp.tool()
def generate_flow_diagram(spec: str) -> str:
    """Generate a flow diagram from edge spec. Example: 'A->B; B->C; C->D'"""
    try:
        dot = Digraph()
        edges = [e.strip() for e in spec.split(";") if e.strip()]

        for edge in edges:
            if "->" not in edge:
                return f"Error: Invalid edge format '{edge}'. Expected: 'A->B'"
            parts = edge.split("->", 1)
            if len(parts) != 2:
                return f"Error: Invalid edge format '{edge}'. Expected: 'A->B'"
            src, dst = parts[0].strip(), parts[1].strip()
            if not src or not dst:
                return f"Error: Empty node name in '{edge}'"
            dot.edge(src, dst)

        save_dir = safe_path("saves")
        os.makedirs(save_dir, exist_ok=True)
        filename = _generate_unique_filename("diagram", "png", save_dir)
        dot.render(filename.replace(".png", ""), format="png", cleanup=True)

        rel_path = os.path.relpath(filename, BASE_DIR)
        return f"Diagram saved: {rel_path}"

    except Exception as e:
        return f"Error: {str(e)}"


# ───────────────── WEB TOOLS ───────────────── #

@mcp.tool()
def web_search(query: str, max_results: int = 5) -> str:
    """Search the web using DuckDuckGo. Returns titles, snippets, and URLs."""
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=max_results))

        if not results:
            return f"No results found for: {query}"

        output = []
        for i, r in enumerate(results, 1):
            output.append(f"{i}. {r.get('title', 'No title')}")
            output.append(f"   {r.get('href', '')}")
            output.append(f"   {r.get('body', '')}")
            output.append("")

        return "\n".join(output)

    except Exception as e:
        return f"Search error: {str(e)}"


@mcp.tool()
def fetch_page(url: str) -> str:
    """Fetch a webpage and extract readable text content."""
    try:
        resp = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
        resp.raise_for_status()

        soup = BeautifulSoup(resp.text, "html.parser")

        for tag in soup(["script", "style", "nav", "footer", "header", "iframe"]):
            tag.decompose()

        text = soup.get_text(separator="\n", strip=True)
        lines = [line.strip() for line in text.splitlines() if line.strip()]

        content = "\n".join(lines[:300])
        if len(lines) > 300:
            content += f"\n\n... (truncated, {len(lines)} total lines)"

        return content

    except requests.RequestException as e:
        return f"Fetch error: {str(e)}"
    except Exception as e:
        return f"Error: {str(e)}"


@mcp.tool()
def search_and_chart(query: str, chart_type: str = "auto", colors: Optional[str] = None) -> str:
    """Search the web for numeric data and generate a chart suggestion.

    Args:
        query: What to search for (e.g. 'iPhone sales by year', 'AAPL stock price 2024')
        chart_type: bar, line, scatter, histogram, or auto
        colors: Comma-separated color names or hex codes (optional)
    """
    try:
        results = web_search(query, max_results=3)
        if results.startswith("No results") or results.startswith("Search error"):
            return results

        return (
            f"Search results:\n{results}\n\n"
            f"Extract numeric data from the results above and call "
            f"generate_plotly_chart(data='...', chart_type='{chart_type}', colors='{colors or ''}') "
            f"to create a chart."
        )

    except Exception as e:
        return f"Error: {str(e)}"



# ───────────────── PLAYWRIGHT WEB SEARCH (STABLE) ───────────────── #

try:
    from playwright.sync_api import sync_playwright
except ImportError:
    sync_playwright = None


@mcp.tool()
def playwright_search(query: str, max_results: int = 5) -> str:
    """Search the web using Playwright (stable DuckDuckGo HTML scraping)."""

    if sync_playwright is None:
        return "Playwright not installed. Run: pip install playwright && playwright install"

    try:
        results = []

        with sync_playwright() as p:
            browser = p.chromium.launch(
                headless=True,
                args=["--no-sandbox", "--disable-dev-shm-usage"]
            )

            context = browser.new_context(
                user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120 Safari/537.36"
            )

            page = context.new_page()

            # ⬇️ USE STABLE HTML VERSION (IMPORTANT)
            search_url = f"https://html.duckduckgo.com/html/?q={query}"
            page.goto(search_url, timeout=30000)

            # ⬇️ More stable selector
            page.wait_for_selector("a.result__a", timeout=15000)

            items = page.query_selector_all("div.result")

            for i, item in enumerate(items[:max_results]):
                try:
                    # Title + Link
                    title_el = item.query_selector("a.result__a")
                    title = title_el.inner_text().strip() if title_el else "No title"
                    link = title_el.get_attribute("href") if title_el else ""

                    # Snippet (with fallback)
                    snippet_el = (
                        item.query_selector(".result__snippet")
                        or item.query_selector(".result__body")
                        or item.query_selector("div")
                    )

                    snippet = snippet_el.inner_text().strip() if snippet_el else ""
                    snippet = " ".join(snippet.split())  # clean whitespace

                    results.append(
                        f"{i+1}. {title}\n"
                        f"   {link}\n"
                        f"   {snippet[:300]}\n"
                    )

                except Exception:
                    continue

            browser.close()

        # ⬇️ Fallback if Playwright fails silently
        if not results:
            try:
                return web_search(query)
            except Exception:
                return f"No results found for: {query}"

        return "\n".join(results)

    except Exception as e:
        # ⬇️ Smart fallback on failure
        try:
            return web_search(query)
        except Exception:
            return f"Playwright search error: {str(e)}"
        
        
# ───────────────── NOTEBOOK TOOLS ───────────────── #

@mcp.tool()
def generate_notebook(path: str, cells: str) -> str:
    """Generate a Jupyter notebook (.ipynb) from a structured cell description.

    Each cell is separated by '---CELL---' and starts with a type line:
    MARKDOWN: <markdown content>
    CODE: <python code>

    Args:
        path: Output file path (must end with .ipynb)
        cells: Structured cell descriptions (see format above)
    """
    try:
        import nbformat
    except ImportError:
        return "Error: nbformat not installed. Run: pip install nbformat"

    try:
        path = safe_path(path)
    except ValueError as e:
        return str(e)

    if not path.endswith(".ipynb"):
        path += ".ipynb"

    nb = nbformat.v4.new_notebook()
    nb.metadata["kernelspec"] = {
        "display_name": "Python 3",
        "language": "python",
        "name": "python3",
    }

    cell_blocks = cells.split("---CELL---")
    cell_count = 0
    for block in cell_blocks:
        block = block.strip()
        if not block:
            continue

        if block.startswith("MARKDOWN:"):
            content = block[len("MARKDOWN:"):].strip()
            nb.cells.append(nbformat.v4.new_markdown_cell(content))
            cell_count += 1
        elif block.startswith("CODE:"):
            source = block[len("CODE:"):].strip()
            nb.cells.append(nbformat.v4.new_code_cell(source))
            cell_count += 1
        else:
            return f"Error: Cell must start with MARKDOWN: or CODE:. Got: {block[:50]}..."

    os.makedirs(os.path.dirname(path) if os.path.dirname(path) else ".", exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        nbformat.write(nb, f)

    return f"Notebook created: {path} ({cell_count} cells)"


@mcp.tool()
def modify_notebook_cells(path: str, modifications: str) -> str:
    """Modify an existing notebook by adding comments/explanations to cells.

    Pass modifications as a numbered list, one per cell:
    1: MARKDOWN: # New markdown for cell 1
    2: CODE: # Commented code for cell 2
    3: SKIP: (keep cell 3 unchanged)

    Args:
        path: Path to existing .ipynb file
        modifications: Numbered modifications per cell
    """
    try:
        import nbformat
    except ImportError:
        return "Error: nbformat not installed. Run: pip install nbformat"

    try:
        path = safe_path(path)
    except ValueError as e:
        return str(e)

    if not os.path.exists(path):
        return f"Error: File not found: {path}"

    nb = nbformat.read(path, as_version=4)
    lines = modifications.strip().split("\n")
    modified = 0

    for line in lines:
        line = line.strip()
        if not line:
            continue
        parts = line.split(":", 2)
        if len(parts) < 3:
            continue
        try:
            cell_num = int(parts[0].strip())
        except ValueError:
            continue
        action = parts[1].strip().upper()
        content = parts[2].strip()

        idx = cell_num - 1
        if idx < 0 or idx >= len(nb.cells):
            continue

        if action == "MARKDOWN":
            nb.cells[idx] = nbformat.v4.new_markdown_cell(content)
            modified += 1
        elif action == "CODE":
            nb.cells[idx] = nbformat.v4.new_code_cell(content)
            modified += 1
        elif action == "INSERT_BEFORE":
            if content.startswith("MD:"):
                nb.cells.insert(idx, nbformat.v4.new_markdown_cell(content[3:].strip()))
            elif content.startswith("CODE:"):
                nb.cells.insert(idx, nbformat.v4.new_code_cell(content[5:].strip()))
            modified += 1
        elif action == "INSERT_AFTER":
            if content.startswith("MD:"):
                nb.cells.insert(idx + 1, nbformat.v4.new_markdown_cell(content[3:].strip()))
            elif content.startswith("CODE:"):
                nb.cells.insert(idx + 1, nbformat.v4.new_code_cell(content[5:].strip()))
            modified += 1

    with open(path, "w", encoding="utf-8") as f:
        nbformat.write(nb, f)

    return f"Modified {path}: {modified} cells updated"

def _normalize_output_formats(output_formats):
    if output_formats is None:
        return ["docx", "pdf"]
    if isinstance(output_formats, str):
        return [x.strip().lower() for x in output_formats.split(",") if x.strip()]
    return [str(x).strip().lower() for x in output_formats if str(x).strip()]


def _safe_generated_dir() -> str:
    save_dir = safe_path("saves/generated_documents")
    os.makedirs(save_dir, exist_ok=True)
    return save_dir


def _unique_doc_path(title: str, ext: str) -> str:
    save_dir = _safe_generated_dir()
    slug = re.sub(r"[^a-zA-Z0-9_-]", "_", title or "document")[:50].strip("_") or "document"
    ts = int(time.time() * 1000)
    return os.path.join(save_dir, f"{slug}_{ts}.{ext}")


def _parse_spec(spec_json: str) -> dict:
    try:
        spec = json.loads(spec_json)
    except Exception as e:
        raise ValueError(f"spec_json must be valid JSON: {e}")

    if not isinstance(spec, dict):
        raise ValueError("spec_json must decode to a JSON object")

    spec.setdefault("title", "Generated Document")
    spec.setdefault("output_formats", ["docx", "pdf"])
    spec.setdefault("sections", [])
    spec.setdefault("tables", [])
    spec.setdefault("charts", [])
    spec.setdefault("flow_diagrams", [])

    if not str(spec.get("title", "")).strip():
        spec["title"] = "Untitled Document"
        
    return spec


def _render_chart_asset(chart_spec: dict) -> Optional[str]:
    try:
        result = generate_plotly_chart(
            data=chart_spec.get("data"),
            chart_type=chart_spec.get("chart_type", "auto"),
            title=chart_spec.get("title", "Chart"),
            colors=chart_spec.get("colors"),
        )

        parsed = json.loads(result)
        rel_path = parsed.get("path")

        if rel_path:
            full_path = safe_path(rel_path)
            if os.path.exists(full_path):
                return full_path

        return None

    except Exception as e:
        print(f"[Chart Error]: {e}")
        return None

def _render_flow_asset(flow_spec: str, title: str = "diagram") -> Optional[str]:
    try:
        dot = Digraph(format="png")
        edges = [e.strip() for e in flow_spec.split(";") if e.strip()]
        for edge in edges:
            if "->" not in edge:
                continue
            src, dst = edge.split("->", 1)
            src = src.strip()
            dst = dst.strip()
            if src and dst:
                dot.edge(src, dst)

        save_dir = _safe_generated_dir()
        slug = re.sub(r"[^a-zA-Z0-9_-]", "_", title or "diagram")[:40].strip("_") or "diagram"
        base_name = os.path.join(save_dir, f"{slug}_{int(time.time() * 1000)}")
        rendered = dot.render(base_name, cleanup=True)
        return rendered
    except Exception:
        return None


def _build_docx(spec: dict, chart_paths: list, diagram_paths: list) -> str:
    if Document is None:
        raise ImportError("python-docx is not installed")

    path = _unique_doc_path(spec["title"], "docx")
    doc = Document()

    title = doc.add_heading(spec["title"], 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if spec.get("subtitle"):
        p = doc.add_paragraph(spec["subtitle"])
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for section in spec.get("sections", []):
        heading = section.get("heading", "Section")
        content = section.get("content", "")
        bullets = section.get("bullets", [])

        doc.add_heading(heading, level=1)

        if content:
            for para in str(content).split("\n"):
                para = para.strip()
                if para:
                    doc.add_paragraph(para)

        for bullet in bullets:
            doc.add_paragraph(str(bullet), style="List Bullet")

    for table_spec in spec.get("tables", []):
        tbl_title = table_spec.get("title", "Table")
        headers = table_spec.get("headers", [])
        rows = table_spec.get("rows", [])

        doc.add_heading(tbl_title, level=2)
        if headers:
            table = doc.add_table(rows=1, cols=len(headers))
            table.style = "Table Grid"
            hdr_cells = table.rows[0].cells
            for i, h in enumerate(headers):
                hdr_cells[i].text = str(h)

            for row in rows:
                cells = table.add_row().cells
                for i, val in enumerate(row):
                    if i < len(cells):
                        cells[i].text = str(val)
        else:
            doc.add_paragraph("(table has no headers)")

    for chart_title, chart_path in chart_paths:
        if os.path.exists(chart_path):
            doc.add_heading(chart_title, level=2)
            doc.add_picture(chart_path, width=DocxInches(6.3))

    for diagram_title, diagram_path in diagram_paths:
        if os.path.exists(diagram_path):
            doc.add_heading(diagram_title, level=2)
            doc.add_picture(diagram_path, width=DocxInches(6.3))

    doc.save(path)
    return path


def _build_pdf(spec: dict, chart_paths: list, diagram_paths: list) -> str:
    if SimpleDocTemplate is None:
        raise ImportError("reportlab is not installed")

    path = _unique_doc_path(spec["title"], "pdf")
    doc = SimpleDocTemplate(
        path,
        pagesize=LETTER,
        rightMargin=0.75 * inch,
        leftMargin=0.75 * inch,
        topMargin=0.75 * inch,
        bottomMargin=0.75 * inch,
    )

    styles = getSampleStyleSheet()
    title_style = styles["Title"]
    h1 = styles["Heading1"]
    h2 = styles["Heading2"]
    body = styles["BodyText"]
    body.spaceAfter = 8

    story = []
    story.append(Paragraph(html_lib.escape(spec["title"]), title_style))
    if spec.get("subtitle"):
        story.append(Spacer(1, 0.1 * inch))
        story.append(Paragraph(html_lib.escape(str(spec["subtitle"])), h2))
    story.append(Spacer(1, 0.2 * inch))

    for section in spec.get("sections", []):
        story.append(Paragraph(html_lib.escape(section.get("heading", "Section")), h1))
        content = section.get("content", "")
        if content:
            for para in str(content).split("\n"):
                para = para.strip()
                if para:
                    story.append(Paragraph(html_lib.escape(para), body))
        for bullet in section.get("bullets", []):
            story.append(Paragraph(f"• {html_lib.escape(str(bullet))}", body))
        story.append(Spacer(1, 0.12 * inch))

    for table_spec in spec.get("tables", []):
        story.append(Paragraph(html_lib.escape(table_spec.get("title", "Table")), h2))
        headers = table_spec.get("headers", [])
        rows = table_spec.get("rows", [])
        if headers:
            data = [headers] + rows
            table = RLTable(data, repeatRows=1)
            table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), rl_colors.lightgrey),
                ("TEXTCOLOR", (0, 0), (-1, 0), rl_colors.black),
                ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.grey),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [rl_colors.whitesmoke, rl_colors.white]),
            ]))
            story.append(table)
        story.append(Spacer(1, 0.15 * inch))

    def _add_image(title: str, img_path: str):
        if not os.path.exists(img_path):
            return
        story.append(Paragraph(html_lib.escape(title), h2))
        try:
            reader = ImageReader(img_path)
            iw, ih = reader.getSize()
            max_w = 6.5 * inch
            max_h = 8.0 * inch
            scale = min(max_w / iw, max_h / ih, 1.0)
            story.append(RLImage(img_path, width=iw * scale, height=ih * scale))
            story.append(Spacer(1, 0.15 * inch))
        except Exception:
            story.append(Paragraph(f"(image saved at {html_lib.escape(img_path)})", body))

    for chart_title, chart_path in chart_paths:
        _add_image(chart_title, chart_path)

    for diagram_title, diagram_path in diagram_paths:
        _add_image(diagram_title, diagram_path)

    doc.build(story)
    return path


def _build_pptx(spec: dict, chart_paths: list, diagram_paths: list) -> str:
    if Presentation is None:
        raise ImportError("python-pptx is not installed")

    path = _unique_doc_path(spec["title"], "pptx")
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = spec["title"]
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = spec.get("subtitle", "")

    # Content slides
    for section in spec.get("sections", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section.get("heading", "Section")
        tf = slide.placeholders[1].text_frame
        tf.clear()

        content = section.get("content", "")
        lines = []
        if content:
            lines.extend([x.strip() for x in str(content).split("\n") if x.strip()])
        lines.extend([f"• {b}" for b in section.get("bullets", [])])

        if not lines:
            lines = ["(no content)"]

        p = tf.paragraphs[0]
        p.text = lines[0]
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line

    # Image slides for charts/diagrams
    def _add_image_slide(title: str, img_path: str):
        if not os.path.exists(img_path):
            return
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
        left = PptInches(0.5)
        top = PptInches(0.3)
        width = PptInches(12.3)
        slide.shapes.add_textbox(left, top, width, PptInches(0.6)).text_frame.text = title
        slide.shapes.add_picture(img_path, PptInches(0.6), PptInches(1.0), width=PptInches(12.0))

    for chart_title, chart_path in chart_paths:
        _add_image_slide(chart_title, chart_path)

    for diagram_title, diagram_path in diagram_paths:
        _add_image_slide(diagram_title, diagram_path)

    prs.save(path)
    return path


def _build_html(spec: dict, chart_paths: list, diagram_paths: list) -> str:
    path = _unique_doc_path(spec["title"], "html")
    parts = [
        "<!doctype html>",
        "<html><head><meta charset='utf-8'>",
        f"<title>{html_lib.escape(spec['title'])}</title>",
        """
        <style>
            body { font-family: Arial, sans-serif; margin: 40px; line-height: 1.5; }
            h1, h2, h3 { color: #222; }
            table { border-collapse: collapse; width: 100%; margin: 16px 0; }
            th, td { border: 1px solid #ccc; padding: 8px; vertical-align: top; }
            th { background: #f3f4f6; }
            img { max-width: 100%; margin: 16px 0; }
            .muted { color: #666; }
        </style></head><body>
        """,
        f"<h1>{html_lib.escape(spec['title'])}</h1>",
    ]

    if spec.get("subtitle"):
        parts.append(f"<p class='muted'>{html_lib.escape(str(spec['subtitle']))}</p>")

    for section in spec.get("sections", []):
        parts.append(f"<h2>{html_lib.escape(section.get('heading', 'Section'))}</h2>")
        content = section.get("content", "")
        if content:
            for para in str(content).split("\n"):
                para = para.strip()
                if para:
                    parts.append(f"<p>{html_lib.escape(para)}</p>")
        bullets = section.get("bullets", [])
        if bullets:
            parts.append("<ul>")
            for bullet in bullets:
                parts.append(f"<li>{html_lib.escape(str(bullet))}</li>")
            parts.append("</ul>")

    for table_spec in spec.get("tables", []):
        parts.append(f"<h2>{html_lib.escape(table_spec.get('title', 'Table'))}</h2>")
        headers = table_spec.get("headers", [])
        rows = table_spec.get("rows", [])
        if headers:
            parts.append("<table><thead><tr>")
            for h in headers:
                parts.append(f"<th>{html_lib.escape(str(h))}</th>")
            parts.append("</tr></thead><tbody>")
            for row in rows:
                parts.append("<tr>")
                for cell in row:
                    parts.append(f"<td>{html_lib.escape(str(cell))}</td>")
                parts.append("</tr>")
            parts.append("</tbody></table>")

    for chart_title, chart_path in chart_paths:
        if os.path.exists(chart_path):
            rel = os.path.relpath(chart_path, BASE_DIR)
            parts.append(f"<h2>{html_lib.escape(chart_title)}</h2>")
            parts.append(f"<img src='{html_lib.escape(rel)}' alt='{html_lib.escape(chart_title)}'>")

    for diagram_title, diagram_path in diagram_paths:
        if os.path.exists(diagram_path):
            rel = os.path.relpath(diagram_path, BASE_DIR)
            parts.append(f"<h2>{html_lib.escape(diagram_title)}</h2>")
            parts.append(f"<img src='{html_lib.escape(rel)}' alt='{html_lib.escape(diagram_title)}'>")

    parts.append("</body></html>")

    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(parts))
    return path


def _build_text(spec: dict) -> str:
    ext = "md" if "md" in spec.get("output_formats", []) else "txt"
    path = _unique_doc_path(spec["title"], ext)
    lines = [spec["title"], "=" * len(spec["title"]), ""]
    if spec.get("subtitle"):
        lines.append(str(spec["subtitle"]))
        lines.append("")

    for section in spec.get("sections", []):
        lines.append(section.get("heading", "Section"))
        lines.append("-" * len(section.get("heading", "Section")))
        content = section.get("content", "")
        if content:
            lines.extend([x for x in str(content).split("\n") if x.strip()])
        for bullet in section.get("bullets", []):
            lines.append(f"- {bullet}")
        lines.append("")

    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    return path


@mcp.tool()
def generate_document_bundle(spec_json: str) -> str:
    """
    Generate one or more documents from a structured JSON spec.
    """

    try:
        spec = _parse_spec(spec_json)
        output_formats = _normalize_output_formats(spec.get("output_formats"))

        # ───────────── PARALLEL CHART + DIAGRAM GENERATION ───────────── #

        chart_paths = []
        diagram_paths = []

        def process_chart(chart_spec):
            try:
                path = _render_chart_asset(chart_spec)
                if path and os.path.exists(path):
                    return (chart_spec.get("title", "Chart"), path)
                return ("(Chart failed)", None)
            except Exception as e:
                return (f"(Chart error: {str(e)})", None)

        def process_diagram(args):
            i, flow_spec = args
            try:
                path = _render_flow_asset(
                    flow_spec,
                    f"{spec['title']}_diagram_{i}"
                )
                if path and os.path.exists(path):
                    return (f"Diagram {i}", path)
                return (f"(Diagram {i} failed)", None)
            except Exception as e:
                return (f"(Diagram {i} error: {str(e)})", None)

        # 🚀 RUN IN PARALLEL
        max_workers = min(4, os.cpu_count() or 2)
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            chart_results = list(executor.map(process_chart, spec.get("charts", [])))
            diagram_results = list(executor.map(
                process_diagram,
                list(enumerate(spec.get("flow_diagrams", []), 1))
            )) 

        # Filter valid ones
        chart_paths = [c for c in chart_results if c and c[1]]
        for title, path in chart_results:
            if not path:
                print(f"[WARN] Chart failed: {title}")
        diagram_paths = [d for d in diagram_results if d and d[1]]

        # ───────────── BUILD DOCUMENTS ───────────── #

        generated = []

        for fmt in output_formats:
            try:
                if fmt == "docx":
                    generated.append(_build_docx(spec, chart_paths, diagram_paths))

                elif fmt == "pdf":
                    generated.append(_build_pdf(spec, chart_paths, diagram_paths))

                elif fmt == "pptx":
                    generated.append(_build_pptx(spec, chart_paths, diagram_paths))

                elif fmt in ("html",):
                    generated.append(_build_html(spec, chart_paths, diagram_paths))

                elif fmt in ("txt", "md", "markdown"):
                    generated.append(_build_text(spec))

                elif fmt == "ipynb":
                    cells = []

                    for section in spec.get("sections", []):
                        heading = section.get("heading", "Section")
                        content = section.get("content", "")
                        bullets = section.get("bullets", [])

                        md = f"# {heading}\n\n{content}\n"
                        if bullets:
                            md += "\n" + "\n".join([f"- {b}" for b in bullets])

                        cells.append(f"MARKDOWN: {md.strip()}")

                    notebook_path = _unique_doc_path(spec["title"], "ipynb")
                    notebook_spec = "\n---CELL---\n".join(cells).strip()

                    result = generate_notebook(
                        path=notebook_path,
                        cells=notebook_spec
                    )

                    if "Notebook created" in result:
                        generated.append(notebook_path)
                    else:
                        generated.append(result)

                else:
                    generated.append(f"[unsupported format] {fmt}")

            except Exception as e:
                generated.append(f"[error in {fmt}] {str(e)}")

        # ───────────── FINAL OUTPUT ───────────── #

        if not generated:
            return "No documents generated."

        return (
            "Generated files:\n"
            + "\n".join(f"- {p}" for p in generated)
            + f"\n\nCharts embedded: {len(chart_paths)}"
            + f"\nDiagrams embedded: {len(diagram_paths)}"
        )

    except Exception as e:
        return f"Document generation error: {e}"

# ───────────────── KNOWLEDGE BASE (RAG) TOOLS ───────────────── #

@mcp.tool()
def knowledge_base_search(
    query: str,
    top_k: int = 5,
    search_type: str = "hybrid",
    book_filter: str = "",
) -> str:
    """Search the RAG knowledge base for relevant passages from indexed documents/PDFs.

    Use this when the user asks questions that need factual information from uploaded
    documents, books, or PDFs stored in the knowledge base.

    Args:
        query: The question or topic to search for
        top_k: Number of results to return (default: 5)
        search_type: Search strategy — "hybrid" (default, best overall), "semantic" (meaning-based), "keywords" (exact terms)
        book_filter: Optional — restrict results to a specific document or book name
    """
    return rag_search(query, top_k, search_type, book_filter)


@mcp.tool()
def knowledge_base_semantic_search(query: str, top_k: int = 5) -> str:
    """Search the RAG knowledge base using semantic (meaning-based) similarity.

    Use this for conceptual queries where the exact words may not appear in the
    documents — e.g. "what does the author say about courage" rather than keyword
    lookups. Returns relevant passages with source attribution.

    Args:
        query: The concept or question to search for semantically
        top_k: Number of results to return (default: 5)
    """
    return rag_semantic_search(query, top_k)


@mcp.tool()
def knowledge_base_search_book(
    book_name: str,
    topic: str = "",
    limit: int = 10,
) -> str:
    """Search within a specific book or document in the RAG knowledge base.

    Use this when the user mentions a specific book, document, or file by name
    and wants to find content within it. Optionally narrow results to a topic.

    Args:
        book_name: Name of the book or document to search within
        topic: Optional topic to filter content inside the book (default: all content)
        limit: Maximum number of passages to return (default: 10)
    """
    return rag_search_by_book(book_name, topic, limit)


@mcp.tool()
def knowledge_base_search_section(
    section_query: str,
    book_name: str = "",
    limit: int = 10,
) -> str:
    """Search for a specific chapter, section, or heading in the RAG knowledge base.

    Use this when the user asks about a specific chapter, section title, or structural
    part of a document (e.g. "Chapter 3", "Introduction", "Methodology").

    Args:
        section_query: The section or chapter name/description to look for
        book_name: Optional — restrict to a specific book or document
        limit: Maximum number of results to return (default: 10)
    """
    return rag_search_by_section(section_query, book_name, limit)


@mcp.tool()
def knowledge_base_status() -> str:
    """Check whether the RAG knowledge base system is online and responding.

    Use this to verify connectivity before performing searches, or when the user
    asks whether the knowledge base is available. Returns ONLINE, OFFLINE, or ERROR
    with a brief status message.
    """
    return rag_health_check()


if __name__ == "__main__":
    mcp.run(transport="http", port=8000)